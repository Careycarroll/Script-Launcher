"""Regenerate test fixtures. Run once after fresh clone."""
from pathlib import Path
import fitz  # pymupdf
from PIL import Image

HERE = Path(__file__).parent

def make_pdf(path: Path, pages: list[str]):
    doc = fitz.open()
    for text in pages:
        page = doc.new_page()
        page.insert_text((72, 72), text, fontsize=14)
    doc.save(path)
    doc.close()

def make_pptx(path: Path):
    # Generate a minimal PPTX via python-pptx if available, else skip
    try:
        from pptx import Presentation
        prs = Presentation()
        for i, title in enumerate(["Slide One", "Slide Two", "Slide Three"]):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
        prs.save(path)
    except ImportError:
        print(f"SKIP {path.name} — python-pptx not installed")

def make_image(path: Path, color: tuple[int, int, int]):
    img = Image.new("RGB", (10, 10), color)
    img.save(path)

if __name__ == "__main__":
    make_pdf(HERE / "sample.pdf", ["Hello from page 1.", "Page two content."])
    make_pdf(HERE / "sample2.pdf", ["Second document, single page."])
    make_pptx(HERE / "sample.pptx")
    make_image(HERE / "red.png", (255, 0, 0))
    make_image(HERE / "blue.png", (0, 0, 255))
    print("done")
